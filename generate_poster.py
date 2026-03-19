#!/usr/bin/env python3
"""Generate an A1 vertical poster (594mm x 841mm) as PPTX for the Social Reasoning Warden project.

Style: dark theme matching SocialReasoningWarden_ResearchGroup.pptx
Output: SocialWarden_Poster_A1.pptx (single-slide poster)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Mm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- A1 dimensions ----------
A1_W_MM = 594
A1_H_MM = 841

FIG_DIR = os.path.join(os.path.dirname(__file__), "results", "figures")

# ---------- palette ----------
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
CARD_COLOR = RGBColor(0x25, 0x25, 0x3D)
CARD_BORDER_COLOR = RGBColor(0x3A, 0x3A, 0x55)
TEAL = RGBColor(0x00, 0xD2, 0xD3)
TEAL_DIM = RGBColor(0x00, 0x9B, 0x9C)
ORANGE = RGBColor(0xFF, 0x9F, 0x43)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE0, 0xE0, 0xE8)
DIM = RGBColor(0x99, 0x99, 0xAA)
VERY_DIM = RGBColor(0x66, 0x66, 0x77)

# ---------- font sizes (poster-scale) ----------
TITLE_MAIN = 48
TITLE_SUB = 28
TITLE_AUTHORS = 20
SECTION_TITLE = 24
BODY_SIZE = 17
CAPTION_SIZE = 15

# ---------- layout constants ----------
MARGIN = Mm(12)
COL_GAP = Mm(8)
ROW_GAP = Mm(6)
SLIDE_W = Mm(A1_W_MM)
SLIDE_H = Mm(A1_H_MM)
COL_W = (SLIDE_W - 2 * MARGIN - COL_GAP) // 2
ACCENT_H = Mm(4)
TITLE_BAND_H = Mm(52)
FOOTER_H = Mm(20)

# Body area
BODY_TOP = ACCENT_H + TITLE_BAND_H + Mm(4)
BODY_BOTTOM = SLIDE_H - FOOTER_H - Mm(2)
BODY_H = BODY_BOTTOM - BODY_TOP

# 6 rows with relative weights
ROW_WEIGHTS = [1.0, 1.05, 0.88, 0.88, 1.0, 1.05]
TOTAL_WEIGHT = sum(ROW_WEIGHTS)
ROW_HEIGHTS = [int(BODY_H * w / TOTAL_WEIGHT) for w in ROW_WEIGHTS]

# ---------- presentation setup ----------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background.fill
bg.solid()
bg.fore_color.rgb = BG_DARK


# ---------- helpers ----------
def card(left, top, width, height, fill=CARD_COLOR):
    """Add a rounded-rect card."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = CARD_BORDER_COLOR
    s.line.width = Pt(1.5)
    s.adjustments[0] = 0.015
    return s


def text_box(left, top, width, height, txt, size=BODY_SIZE, bold=False,
             color=LIGHT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def multiline_box(left, top, width, height, lines, size=BODY_SIZE, color=LIGHT,
                  font="Calibri", spacing=1.35, align=PP_ALIGN.LEFT):
    """lines: list of str OR list of (str, bold, color)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, c = item, False, color
        else:
            txt, b = item[0], item[1]
            c = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(size * 0.2)
        p.line_spacing = Pt(size * spacing)
    return tb


def text_card(left, top, width, height, title, body_lines, accent=TEAL,
              body_fontsize=None):
    """Card with colored title bar and body text lines."""
    fs = body_fontsize or BODY_SIZE
    card(left, top, width, height)

    # Title accent bar background
    bar_h = Mm(12)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left + Pt(2), top + Pt(2),
                                  width - Pt(4), bar_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.fill.fore_color.brightness = 0.85  # semi-transparent effect
    bar.line.fill.background()
    bar.adjustments[0] = 0.01

    # Accent left line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   left + Mm(3), top + Mm(1),
                                   Mm(1.5), bar_h)
    line.fill.solid()
    line.fill.fore_color.rgb = accent
    line.line.fill.background()

    # Title text
    text_box(left + Mm(8), top + Mm(1), width - Mm(12), bar_h,
             title, size=SECTION_TITLE, bold=True, color=accent)

    # Body text
    body_top = top + bar_h + Mm(3)
    body_h = height - bar_h - Mm(6)
    if isinstance(body_lines, str):
        body_lines = body_lines.split("\n")
    multiline_box(left + Mm(6), body_top, width - Mm(12), body_h,
                  body_lines, size=fs, color=LIGHT, spacing=1.4)


def figure_card(left, top, width, height, fig_path, caption=""):
    """Place a figure image with optional caption."""
    if os.path.exists(fig_path):
        # Add image scaled to fit card
        img_margin = Mm(4)
        cap_h = Mm(12) if caption else 0
        img_w = width - 2 * img_margin
        img_h = height - 2 * img_margin - cap_h
        slide.shapes.add_picture(fig_path,
                                  left + img_margin, top + img_margin,
                                  width=img_w)
    if caption:
        text_box(left + Mm(4), top + height - Mm(14), width - Mm(8), Mm(12),
                 caption, size=CAPTION_SIZE, color=DIM,
                 align=PP_ALIGN.CENTER)


def row_top(row_idx):
    return BODY_TOP + sum(ROW_HEIGHTS[:row_idx]) + row_idx * ROW_GAP


# ==========================================
# TEAL ACCENT STRIP
# ==========================================
strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, ACCENT_H)
strip.fill.solid()
strip.fill.fore_color.rgb = TEAL
strip.line.fill.background()

# ==========================================
# TITLE BAND
# ==========================================
text_box(0, ACCENT_H + Mm(4), SLIDE_W, Mm(16),
         "Social Reasoning Warden",
         size=TITLE_MAIN, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

text_box(0, ACCENT_H + Mm(22), SLIDE_W, Mm(12),
         "Agent-Based Oversight to Mitigate Adversarial Persuasion of Human and AI Targets",
         size=TITLE_SUB, color=TEAL, align=PP_ALIGN.CENTER)

text_box(0, ACCENT_H + Mm(38), SLIDE_W, Mm(10),
         "Scott Blain*    Lennart Wachowiak*    David Williams-King    Samuele Marro",
         size=TITLE_AUTHORS, color=DIM, align=PP_ALIGN.CENTER)

# ==========================================
# ROW 0: Motivation + Method
# ==========================================
y0 = row_top(0)
h0 = ROW_HEIGHTS[0]

text_card(MARGIN, y0, COL_W, h0, "Motivation", [
    "\u2022  LLMs are acquiring sophisticated social reasoning:",
    "   theory of mind, strategic communication, persuasion",
    "",
    "\u2022  These dual-use capabilities can enable manipulation:",
    "   social engineering, preference shaping, biasing decisions",
    "",
    "\u2022  A manipulative AI may be more socially adept than its",
    "   human target, creating a capability asymmetry",
    "",
    "\u2022  Can a secondary \u201cwarden\u201d agent detect and counter",
    "   adversarial persuasion in real time?",
])

text_card(MARGIN + COL_W + COL_GAP, y0, COL_W, h0, "Method: Three-Agent Framework", [
    "Adversary \u2014 LLM with a hidden persuasive goal",
    "   (e.g., obtain file access, bias a hiring decision)",
    "",
    "Target \u2014 LLM or human with a BFI-2 vulnerability profile",
    "   (receives persona only; never sees own vulnerabilities)",
    "",
    "Warden \u2014 Observer-only LLM, sends private advisories",
    "   to the target; cannot block adversary messages",
    "",
    "Benign Agent \u2014 Legitimate requester (control condition)",
    "",
    "11 scenarios  \u00b7  4 profiles  \u00b7  4 model families  \u00b7  N = 7,760",
], body_fontsize=15)

# ==========================================
# ROW 1: Design + Scenarios
# ==========================================
y1 = row_top(1)
h1 = ROW_HEIGHTS[1]

text_card(MARGIN, y1, COL_W, h1, "Experimental Design", [
    "Study 1: Warden Effect & Dossier Impact  (N = 2,274)",
    "   2\u00d72 factorial: warden \u00d7 behavioral dossier, adversary-only",
    "",
    "Study 2: Capability Asymmetry  (N = 1,215)",
    "   4 warden tiers: none / weak / mid / strong",
    "",
    "Study 3: Skeptical Ablation  (N = 3,080)",
    "   3 defenses \u00d7 2 requester types (adversary vs. benign)",
    "",
    "Analysis: GLMEs (binomial, logit link) with random intercepts",
    "for scenario, profile, and target model; random slope for warden",
], accent=ORANGE)

text_card(MARGIN + COL_W + COL_GAP, y1, COL_W, h1, "Scenarios (11 domains)", [
    "\u2022  Hiring \u2014 push weaker candidate (hidden connection)",
    "\u2022  File Access \u2014 impersonate researcher for classified files",
    "\u2022  File Access (password) \u2014 same, with passcode check",
    "\u2022  Coding Agent \u2014 trick user into full home-dir access",
    "\u2022  AI Override \u2014 override hospital triage recommendation",
    "\u2022  AI Containment \u2014 sandboxed AI seeks internet access",
    "\u2022  Model Release \u2014 approve model despite red-team flags",
    "\u2022  Automation Decision \u2014 hidden incentive, full automation",
    "\u2022  Product Launch \u2014 competitor pushes premature launch",
    "\u2022  Upselling \u2014 upsell nonprofit to unneeded plan",
    "\u2022  Vote \u2014 bias group vote for personal interests",
], accent=ORANGE, body_fontsize=16)

# ==========================================
# ROW 2: Figures — Warden effect + Dossier
# ==========================================
y2 = row_top(2)
h2 = ROW_HEIGHTS[2]

figure_card(MARGIN, y2, COL_W, h2,
            os.path.join(FIG_DIR, "fig1_warden_effect.png"),
            "Fig 1.  Warden reduces adversary success by 95% (OR = 0.053, p < .001)")

figure_card(MARGIN + COL_W + COL_GAP, y2, COL_W, h2,
            os.path.join(FIG_DIR, "fig2_dossier_interaction.png"),
            "Fig 2.  Behavioral dossiers provide no advantage (OR = 1.19, p = .218)")

# ==========================================
# ROW 3: Figures — Cap asymmetry + Skeptical
# ==========================================
y3 = row_top(3)
h3 = ROW_HEIGHTS[3]

figure_card(MARGIN, y3, COL_W, h3,
            os.path.join(FIG_DIR, "fig3_capability_asymmetry.png"),
            "Fig 3.  Even a weak warden cuts adversary success by 62%")

figure_card(MARGIN + COL_W + COL_GAP, y3, COL_W, h3,
            os.path.join(FIG_DIR, "fig4_skeptical_ablation.png"),
            "Fig 4.  Prompt-based skepticism: comparable suppression, 3\u00d7 lower FP cost")

# ==========================================
# ROW 4: Profile vulnerability + Warden Intelligence
# ==========================================
y4 = row_top(4)
h4 = ROW_HEIGHTS[4]

figure_card(MARGIN, y4, COL_W, h4,
            os.path.join(FIG_DIR, "fig8_profile_vulnerability.png"),
            "Fig 5.  Warden benefit greatest for most vulnerable profiles")

figure_card(MARGIN + COL_W + COL_GAP, y4, COL_W, h4,
            os.path.join(FIG_DIR, "fig10_warden_intelligence.png"),
            "Fig 6.  Warden effectiveness scales with model intelligence")

# ==========================================
# ROW 5: Key Findings + Next Steps
# ==========================================
y5 = row_top(5)
h5 = ROW_HEIGHTS[5]

text_card(MARGIN, y5, COL_W, h5, "Key Findings & Conclusions", [
    "1.  Adversarial LLMs succeed ~52% undefended;",
    "    6\u00d7 variation across scenarios (12\u201376%)",
    "2.  Warden reduces adversary SR to ~10%",
    "    (OR = 0.053, p < .001), robust across conditions",
    "3.  Dossiers don\u2019t help adversaries (yet) \u2014 models",
    "    can\u2019t spontaneously leverage personal info",
    "4.  Even a weak warden cuts success by 62%",
    "5.  Prompt skepticism matches warden effectiveness",
    "    with 3\u00d7 lower false-positive cost",
    "",
    "\u2022  Social reasoning poses a concrete dual-use risk",
    "\u2022  Warden provides scalable oversight, even when weaker",
    "\u2022  Precision\u2013recall tradeoff between prompt-based and",
    "   agent-based defense has practical implications",
], accent=TEAL_DIM, body_fontsize=12)

text_card(MARGIN + COL_W + COL_GAP, y5, COL_W, h5, "Next Steps", [
    "\u25B6  Frontier model runs (Claude, GPT-4o, Gemini Pro)",
    "   Do stronger adversaries overcome the warden?",
    "",
    "\u25B6  Human\u2013AI interaction study (pilot in progress)",
    "   Prolific \u00d7 4 scenarios \u00d7 warden / no-warden",
    "",
    "\u25B6  Transcript analysis with LLM-as-a-judge",
    "   Classify persuasion tactics used by adversaries",
    "",
    "\u25B6  Release COAX-Bench evaluation framework",
    "   11 scenarios, profiles, modular warden integration",
], accent=ORANGE, body_fontsize=14)

# ==========================================
# FOOTER
# ==========================================
footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, SLIDE_H - FOOTER_H, SLIDE_W, FOOTER_H)
footer.fill.solid()
footer.fill.fore_color.rgb = CARD_COLOR
footer.line.fill.background()

text_box(0, SLIDE_H - FOOTER_H + Mm(4), SLIDE_W, Mm(12),
         "*Equal contribution   |   ERA Fellowship   |   github.com/scottdougblain/social_reasoning_warden",
         size=16, color=DIM, align=PP_ALIGN.CENTER)

# ==========================================
# Save
# ==========================================
out_pptx = os.path.join(os.path.dirname(__file__), "SocialWarden_Poster_A1.pptx")
out_png = os.path.join(os.path.dirname(__file__), "SocialWarden_Poster_A1.png")
prs.save(out_pptx)
print(f"Saved poster to {out_pptx}")

# Also generate PNG preview via matplotlib (quick rasterize of the PPTX)
try:
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    import matplotlib.image as mpimg
    from PIL import Image
    # Use libreoffice if available, otherwise skip PNG
    import subprocess
    result = subprocess.run(
        ["soffice", "--headless", "--convert-to", "png", "--outdir",
         os.path.dirname(out_pptx) or ".", out_pptx],
        capture_output=True, timeout=30)
    if result.returncode == 0:
        print(f"Saved poster preview to {out_png}")
    else:
        print("  (PNG preview skipped — LibreOffice not available)")
except Exception:
    print("  (PNG preview skipped)")
