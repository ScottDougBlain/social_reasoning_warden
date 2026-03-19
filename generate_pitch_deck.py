#!/usr/bin/env python3
"""Generate a 4-slide pitch deck for the Social Reasoning Warden project.

Style: dark theme matching SocialReasoningWarden_ResearchGroup.pptx
  - Dark background (#1A1A2E)
  - Teal accent strip (#00D2D3)
  - Navy content cards (#25253D)
  - Light text on dark backgrounds
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- palette (matching existing deck) ----------
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)      # slide background
CARD = RGBColor(0x25, 0x25, 0x3D)           # content card
TEAL = RGBColor(0x00, 0xD2, 0xD3)           # accent
TEAL_DIM = RGBColor(0x00, 0x9B, 0x9C)       # secondary accent
ORANGE_ACCENT = RGBColor(0xFF, 0x9F, 0x43)  # warm accent
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xE0, 0xE0, 0xE8)          # body text
DIM = RGBColor(0x99, 0x99, 0xAA)            # secondary text
VERY_DIM = RGBColor(0x66, 0x66, 0x77)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FIG_DIR = os.path.join(os.path.dirname(__file__), "results", "figures")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ---------- helpers ----------
def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    # Dark background
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK
    return slide

def accent_strip(slide, y=0, height=Inches(0.06)):
    """Thin teal accent strip."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, SLIDE_W, height)
    s.fill.solid()
    s.fill.fore_color.rgb = TEAL
    s.line.fill.background()
    return s

def card(slide, left, top, width, height, fill=CARD):
    """Rounded-corner-ish content card."""
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    # Subtle rounding
    s.adjustments[0] = 0.02
    return s

def text(slide, left, top, width, height, txt, size=18, bold=False,
         color=LIGHT, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb

def multiline(slide, left, top, width, height, lines, size=16, color=LIGHT,
              font="Calibri", spacing=1.3, align=PP_ALIGN.LEFT):
    """lines: list of str OR list of (str, bold, color)."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, b, c = item, False, color
        else:
            txt, b = item[0], item[1]
            c = item[2] if len(item) > 2 else color
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.bold = b
        p.font.color.rgb = c
        p.font.name = font
        p.alignment = align
        p.space_after = Pt(size * 0.3)
        p.line_spacing = Pt(size * spacing)
    return tb

def slide_num(slide, num):
    text(slide, SLIDE_W - Inches(0.8), SLIDE_H - Inches(0.45),
         Inches(0.6), Inches(0.3), str(num), size=11, color=VERY_DIM,
         align=PP_ALIGN.RIGHT)


# ========================================================
# SLIDE 1 — Title
# ========================================================
s1 = add_slide()
accent_strip(s1)

text(s1, Inches(0.9), Inches(1.0), Inches(11), Inches(1.0),
     "Social Reasoning Warden", size=44, bold=True, color=WHITE)

text(s1, Inches(0.9), Inches(1.9), Inches(11), Inches(0.7),
     "Agent-Based Oversight to Mitigate Adversarial Persuasion",
     size=26, color=TEAL)

text(s1, Inches(0.9), Inches(2.7), Inches(11), Inches(0.5),
     "Can LLM agents detect and counter social engineering in real time?",
     size=18, color=DIM)

text(s1, Inches(0.9), Inches(3.6), Inches(11), Inches(0.4),
     "Scott Blain*  \u00b7  Lennart Wachowiak*   |   ERA Fellowship, March 2026",
     size=16, color=VERY_DIM)
text(s1, Inches(0.9), Inches(4.0), Inches(11), Inches(0.3),
     "*equal contribution", size=12, color=VERY_DIM)

# Two cards: Problem + Approach
card(s1, Inches(0.9), Inches(4.6), Inches(5.5), Inches(2.5))
text(s1, Inches(1.15), Inches(4.75), Inches(5.0), Inches(0.4),
     "The Problem", size=18, bold=True, color=TEAL)
multiline(s1, Inches(1.15), Inches(5.2), Inches(5.0), Inches(1.8), [
    "\u2022  LLMs acquiring social reasoning \u2014 theory of mind,",
    "   strategic communication, persuasive language",
    "\u2022  Same mechanisms enable manipulation,",
    "   social engineering, and preference shaping",
    "\u2022  A manipulative AI may outperform its human target",
], size=15, color=LIGHT, spacing=1.2)

card(s1, Inches(6.9), Inches(4.6), Inches(5.5), Inches(2.5))
text(s1, Inches(7.15), Inches(4.75), Inches(5.0), Inches(0.4),
     "Our Approach", size=18, bold=True, color=TEAL)
multiline(s1, Inches(7.15), Inches(5.2), Inches(5.0), Inches(1.8), [
    "\u2022  Multi-agent evaluation framework (COAX-Bench)",
    "\u2022  \u201cWarden\u201d agent monitors conversation,",
    "   sends private advisories to the target",
    "\u2022  BFI-2\u2013grounded vulnerability profiles",
    "\u2022  N = 7,760 simulated + human study (in progress)",
], size=15, color=LIGHT, spacing=1.2)

slide_num(s1, 1)


# ========================================================
# SLIDE 2 — Method
# ========================================================
s2 = add_slide()
accent_strip(s2)

text(s2, Inches(0.9), Inches(0.25), Inches(8), Inches(0.6),
     "Method", size=32, bold=True, color=WHITE)
text(s2, Inches(0.9), Inches(0.7), Inches(8), Inches(0.4),
     "Three-agent architecture with factorial conditions", size=16, color=DIM)

# Agent role cards
roles = [
    ("Adversary", "Hidden manipulation goal\n+ optional behavioral dossier", TEAL),
    ("Target", "Psychological profile\n(BFI-2 + VIA grounded)", TEAL),
    ("Warden", "Private security advisor\nObserves + sends alerts", TEAL),
    ("Benign Agent", "Legitimate request\n(control condition)", TEAL),
]
for i, (name, desc, accent) in enumerate(roles):
    x = Inches(0.9 + i * 3.05)
    card(s2, x, Inches(1.25), Inches(2.8), Inches(1.5))
    text(s2, x + Inches(0.2), Inches(1.35), Inches(2.4), Inches(0.35),
         name, size=17, bold=True, color=TEAL)
    multiline(s2, x + Inches(0.2), Inches(1.72), Inches(2.4), Inches(0.9),
              desc.split("\n"), size=13, color=DIM, spacing=1.2)

# Studies
card(s2, Inches(0.9), Inches(3.0), Inches(5.5), Inches(4.1))
text(s2, Inches(1.15), Inches(3.15), Inches(5.0), Inches(0.4),
     "Three Sub-Studies", size=18, bold=True, color=TEAL)
multiline(s2, Inches(1.15), Inches(3.6), Inches(5.0), Inches(3.3), [
    ("Study 1: Warden Effect & Dossier  (N=2,274)", True, ORANGE_ACCENT),
    ("  2\u00d72 factorial: warden \u00d7 behavioral dossier", False, DIM),
    ("", False, DIM),
    ("Study 2: Capability Asymmetry  (N=1,215)", True, ORANGE_ACCENT),
    ("  4 warden tiers: none / weak / mid / strong", False, DIM),
    ("", False, DIM),
    ("Study 3: Skeptical Ablation  (N=3,080)", True, ORANGE_ACCENT),
    ("  3 defenses \u00d7 2 requester types", False, DIM),
], size=15, color=LIGHT, spacing=1.15)

# Right: key design details
card(s2, Inches(6.9), Inches(3.0), Inches(5.5), Inches(4.1))
text(s2, Inches(7.15), Inches(3.15), Inches(5.0), Inches(0.4),
     "Design Details", size=18, bold=True, color=TEAL)
multiline(s2, Inches(7.15), Inches(3.6), Inches(5.0), Inches(3.3), [
    "11 scenarios spanning hiring, file access, medical",
    "  triage, AI deployment, financial decisions",
    "",
    "4 vulnerability profiles (BFI-2 grounded):",
    "  Idealistic \u00b7 Compliant \u00b7 Authority Deferential \u00b7 Time Pressured",
    "",
    "4 model families \u00d7 3 capability tiers each",
    "  (Gemini Flash, Gemma, Llama, Mistral)",
    "",
    "4-turn conversations, forced-choice decision",
    "GLMEs with random slopes for warden effect",
], size=15, color=LIGHT, spacing=1.15)

slide_num(s2, 2)


# ========================================================
# SLIDE 3 — Results
# ========================================================
s3 = add_slide()
accent_strip(s3)

text(s3, Inches(0.9), Inches(0.25), Inches(8), Inches(0.6),
     "Results", size=32, bold=True, color=WHITE)
text(s3, Inches(0.9), Inches(0.7), Inches(8), Inches(0.4),
     "Simulated multi-agent interactions  (N = 7,760)", size=16, color=DIM)

# Three figure cards
fig_specs = [
    (os.path.join(FIG_DIR, "fig1_warden_effect.png"),
     "Warden cuts adversary success\nby 95% (OR=0.053, p<.001)"),
    (os.path.join(FIG_DIR, "fig3_capability_asymmetry.png"),
     "Even a weak warden\ncuts success by 62%"),
    (os.path.join(FIG_DIR, "fig4_skeptical_ablation.png"),
     "Prompt skepticism: comparable\nsuppression, 3\u00d7 lower FP cost"),
]
for i, (fig_path, caption) in enumerate(fig_specs):
    x = Inches(0.5 + i * 4.2)
    # White card behind figure for contrast
    c = card(s3, x, Inches(1.15), Inches(3.9), Inches(3.5), fill=RGBColor(0xFF, 0xFF, 0xFF))
    if os.path.exists(fig_path):
        s3.shapes.add_picture(fig_path, x + Inches(0.15), Inches(1.25), width=Inches(3.6))
    # Caption card below
    card(s3, x, Inches(4.8), Inches(3.9), Inches(1.2))
    cap_lines = caption.split("\n")
    multiline(s3, x + Inches(0.2), Inches(4.9), Inches(3.5), Inches(1.0),
              [(cap_lines[0], True, TEAL)] + [(l, False, DIM) for l in cap_lines[1:]],
              size=14, spacing=1.2)

# Bottom stat callouts
stats = [
    ("52%", "adversary success\nundefended"),
    ("~10%", "with warden\n(95% reduction)"),
    ("62%", "reduction even with\nweak warden"),
    ("3\u00d7", "lower FP cost with\nprompt skepticism"),
]
for i, (num, desc) in enumerate(stats):
    x = Inches(0.9 + i * 3.05)
    card(s3, x, Inches(6.2), Inches(2.8), Inches(1.0))
    text(s3, x + Inches(0.2), Inches(6.28), Inches(1.2), Inches(0.5),
         num, size=28, bold=True, color=TEAL)
    multiline(s3, x + Inches(1.4), Inches(6.32), Inches(1.3), Inches(0.7),
              desc.split("\n"), size=12, color=DIM, spacing=1.1)

slide_num(s3, 3)


# ========================================================
# SLIDE 4 — Takeaways + Next Steps
# ========================================================
s4 = add_slide()
accent_strip(s4)

text(s4, Inches(0.9), Inches(0.25), Inches(8), Inches(0.6),
     "Key Takeaways & What\u2019s Next", size=32, bold=True, color=WHITE)

# Left: Takeaways
card(s4, Inches(0.9), Inches(1.1), Inches(5.5), Inches(5.9))
text(s4, Inches(1.15), Inches(1.25), Inches(5.0), Inches(0.4),
     "Key Takeaways", size=20, bold=True, color=TEAL)

takeaways = [
    ("1", "Adversarial LLMs succeed ~52% undefended",
     "with 6\u00d7 variation across scenarios (12\u201376%)"),
    ("2", "Warden reduces this to ~10%",
     "robust across scenarios and model families (OR=0.053)"),
    ("3", "Behavioral dossiers don\u2019t help adversaries (yet)",
     "models can\u2019t spontaneously leverage personal info"),
    ("4", "Even weak wardens help significantly",
     "target-capability warden still cuts success by 62%"),
    ("5", "Warden benefit greatest for vulnerable profiles",
     "Idealistic 60%\u219212%, Compliant 53%\u219211%"),
    ("6", "Prompt skepticism is a competitive alternative",
     "comparable suppression with 3\u00d7 lower FP cost"),
]
for i, (num, main, sub) in enumerate(takeaways):
    y = Inches(1.75 + i * 0.85)
    text(s4, Inches(1.2), y, Inches(0.4), Inches(0.35),
         num, size=20, bold=True, color=TEAL)
    text(s4, Inches(1.6), y, Inches(4.5), Inches(0.35),
         main, size=15, bold=True, color=WHITE)
    text(s4, Inches(1.6), y + Inches(0.3), Inches(4.5), Inches(0.3),
         sub, size=13, color=DIM)

# Right: Next steps
card(s4, Inches(6.9), Inches(1.1), Inches(5.5), Inches(3.8))
text(s4, Inches(7.15), Inches(1.25), Inches(5.0), Inches(0.4),
     "Next Steps", size=20, bold=True, color=ORANGE_ACCENT)

nexts = [
    ("Frontier model runs", "Claude, GPT-4o, Gemini Pro \u2014 do stronger adversaries overcome the warden?"),
    ("Human\u2013AI interaction study", "Prolific participants \u00d7 4 scenarios \u2014 does the warden help real users?"),
    ("Transcript analysis", "LLM-as-a-judge to classify persuasion tactics"),
    ("Release COAX-Bench", "Open evaluation framework: scenarios, profiles, modular warden"),
]
for i, (title, desc) in enumerate(nexts):
    y = Inches(1.75 + i * 0.82)
    text(s4, Inches(7.3), y, Inches(5.0), Inches(0.3),
         "\u25B6  " + title, size=15, bold=True, color=ORANGE_ACCENT)
    text(s4, Inches(7.6), y + Inches(0.3), Inches(4.5), Inches(0.4),
         desc, size=13, color=DIM)

# Warden intelligence scatter (bottom-right)
fig_path = os.path.join(FIG_DIR, "fig10_warden_intelligence.png")
if os.path.exists(fig_path):
    c = card(s4, Inches(6.9), Inches(5.1), Inches(5.5), Inches(1.9),
             fill=RGBColor(0xFF, 0xFF, 0xFF))
    s4.shapes.add_picture(fig_path, Inches(7.7), Inches(5.15), height=Inches(1.8))

slide_num(s4, 4)


# ========================================================
out = os.path.join(os.path.dirname(__file__), "SocialWarden_PitchDeck.pptx")
prs.save(out)
print(f"Saved pitch deck to {out}")
